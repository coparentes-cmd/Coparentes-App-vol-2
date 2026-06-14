#!/usr/bin/env python3
"""Build PowerPoint from docs/instrukcja-nowa-rodzina.md content + screenshots."""

from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets" / "onboarding"
OUTPUT = ROOT / "docs" / "instrukcja-nowa-rodzina.pptx"
OUTPUT_DOWNLOAD = ROOT / "web" / "downloads" / "instrukcja-nowa-rodzina.pptx"

TEAL = RGBColor(0x00, 0xC8, 0x96)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Jak założyć nową rodzinę\nw Coparentes"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(7.6), Inches(1))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "getcoparentes.app"
    sp.font.size = Pt(22)
    sp.font.color.rgb = TEAL
    sp.alignment = PP_ALIGN.CENTER


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    image_name: Optional[str] = None,
    footer: Optional[str] = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = DARK

    text_width = Inches(4.8) if image_name else Inches(9)
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), text_width, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)
        p.level = 0

    if footer:
        fp = tf.add_paragraph()
        fp.text = footer
        fp.font.size = Pt(14)
        fp.font.color.rgb = TEAL
        fp.space_before = Pt(12)

    if image_name:
        img_path = ASSETS / image_name
        if img_path.exists():
            slide.shapes.add_picture(
                str(img_path),
                Inches(5.35),
                Inches(1.0),
                height=Inches(5.8),
            )


def add_text_only_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    add_bullet_slide(prs, title, bullets, image_name=None)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_text_only_slide(
        prs,
        "Wprowadzenie",
        [
            "Instrukcja dla pierwszego rodzica (Rodzic A).",
            "Konto produkcyjne — nie tryb demo.",
            "Hasło: minimum 10 znaków.",
            "Aplikacja: getcoparentes.app",
        ],
    )

    add_bullet_slide(
        prs,
        "Krok 1 — Utwórz konto (Rodzic A)",
        [
            "Wybierz zakładkę „Nowa rodzina”.",
            "Wypełnij: imię i nazwisko, nazwa rodziny, e-mail, hasło.",
            "Kliknij „Utwórz konto”.",
        ],
        "01-nowa-rodzina.png",
        footer="Po rejestracji jesteś Rodzicem A — zarządzasz zaproszeniami.",
    )

    add_bullet_slide(
        prs,
        "Krok 2 — Dodaj pierwsze dziecko (opcjonalnie)",
        [
            "Okno „Dodaj dziecko do przestrzeni” pojawia się po rejestracji.",
            "Podaj imię, datę urodzenia (ważną przy logowaniu dziecka).",
            "Szkoła opcjonalna.",
            "„Dodaj dziecko” lub „Pomiń na razie”.",
        ],
        "02-dodaj-dziecko-sheet.png",
    )

    add_bullet_slide(
        prs,
        "Krok 3A — Zaproś drugiego rodzica (kod)",
        [
            "Ustawienia → ikona zębatki przy awatarze.",
            "„Kod zaproszenia dla drugiego rodzica” → dotknij, aby skopiować.",
            "Wyślij kod partnerowi (SMS, WhatsApp, e-mail).",
            "Ten sam ekran: kod zaproszenia dziecka (krok 6).",
        ],
        "03-ustawienia-kody.png",
    )

    add_bullet_slide(
        prs,
        "Krok 3B — Zaproszenie e-mailem (opcjonalnie)",
        [
            "Ustawienia → „Zaproś drugiego rodzica mailem”.",
            "Partner musi mieć już konto Coparentes.",
            "Przy pierwszym dołączeniu wygodniejszy jest kod (krok 3A).",
        ],
        "04-zaproszenie-email.png",
    )

    add_bullet_slide(
        prs,
        "Krok 4 — Drugi rodzic dołącza",
        [
            "getcoparentes.app → zakładka „Rodzic” (nie „Nowa rodzina”).",
            "Imię, kod od Rodzica A, własny e-mail i hasło.",
            "Kliknij „Dołącz”.",
        ],
        "05-rodzic-dolacz.png",
    )

    add_bullet_slide(
        prs,
        "Krok 5 — Dodaj kolejne dzieci",
        [
            "Ustawienia → „Dodaj dziecko”.",
            "Ten sam formularz co w kroku 2.",
            "Lista dzieci: „Dzieci w rodzinie” w ustawieniach.",
        ],
        "06-dodaj-kolejne-dziecko.png",
    )

    add_bullet_slide(
        prs,
        "Krok 6 — Kod dla dziecka",
        [
            "Ustawienia → skopiuj „Kod zaproszenia dziecka”.",
            "Jeden kod na rodzinę — dziecko rozpoznawane po dacie urodzenia.",
            "Przekaż kod bezpiecznie (ustnie lub wiadomość od rodzica).",
        ],
        "07-kod-dziecka.png",
    )

    add_bullet_slide(
        prs,
        "Krok 7 — Logowanie dziecka",
        [
            "Zakładka „Dziecko” → kod → „Sprawdź kod”.",
            "Data urodzenia (jak w profilu) + hasło (min. 10 znaków).",
            "Pierwsze logowanie: imię. Kolejne: kod + data + hasło.",
            "„Wejdź” → panel: Dzisiaj, Kalendarz, Rodzina, Lista.",
        ],
        "08-dziecko-logowanie.png",
    )

    add_text_only_slide(
        prs,
        "Typowe problemy",
        [
            "Brak profilu dziecka → sprawdź datę urodzenia u Rodzica A.",
            "„Sprawdź kod” → kliknij przed „Wejdź”.",
            "Zły kod → kod dziecka ≠ kod drugiego rodzica.",
            "Hasło za krótkie → min. 10 znaków.",
            "Profil zajęty → dziecko loguje się ponownie tym samym hasłem.",
        ],
    )

    add_text_only_slide(
        prs,
        "Podsumowanie ról",
        [
            "Rodzic A — „Nowa rodzina” → ustawienia: zakłada rodzinę, kody, dzieci.",
            "Rodzic B — „Rodzic”: dołącza kodem od Rodzica A.",
            "Dziecko — „Dziecko”: kod rodziny + data urodzenia + hasło.",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    OUTPUT_DOWNLOAD.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    prs.save(str(OUTPUT_DOWNLOAD))
    print(f"Saved {OUTPUT}")
    print(f"Saved {OUTPUT_DOWNLOAD}")


if __name__ == "__main__":
    main()
